import time
import os
from PIL import ImageGrab, Image
from pynput import keyboard
import numpy as np
from skimage.metrics import structural_similarity as ssim
from pptx import Presentation
from pptx.util import Inches
import cv2
import tkinter as tk
from tkinter import simpledialog
import threading
import sys
import traceback
import shutil

class VideoToPPT:
    def __init__(self):
        self.running = False
        self.images = []
        self.ppt = None
        self.video_region = None
        self.last_image = None
        self.similarity_threshold = 0.85  # Lowered from 0.95 to be less strict on similarity
        self.capture_count = 0
        self.max_capture_time = 60  # Auto-save after 60 seconds regardless
        self.start_time = 0
        self.save_interval = 5  # Save every 5 captures as backup
        self.text_detection_enabled = True  # Enable text-based comparison
        
    def toggle_capture(self):
        """Take a single screenshot instead of continuous capturing."""
        try:
            # Initialize ppt if not already done
            if not self.ppt:
                self.initialize_ppt()
                
            # Detect video region if not already done    
            if not self.video_region:
                self.detect_video_region()
            
            # Take a single screenshot
            print("Taking screenshot...")
            screenshot = ImageGrab.grab(bbox=self.video_region)
            
            # Add screenshot to presentation (no duplicate check)
            self.add_to_presentation(screenshot)
            
            # Save immediately
            self.save_ppt(final=True)
            print("Screenshot captured and saved to PowerPoint.")
            
        except Exception as e:
            print(f"Error capturing screenshot: {e}")
            traceback.print_exc()
            self.emergency_save()
    
    def auto_timeout(self):
        """Automatically stop capture after max_capture_time."""
        try:
            while self.running:
                elapsed = time.time() - self.start_time
                if elapsed > self.max_capture_time:
                    print(f"\nAUTO-TIMEOUT: Capture ran for {elapsed:.1f} seconds")
                    print("Automatically stopping and saving...")
                    self.running = False
                    self.save_ppt(final=True)
                    return
                time.sleep(1)
        except Exception as e:
            print(f"Error in auto_timeout: {e}")
            self.emergency_save()
            
    def initialize_ppt(self):
        """Initialize the PowerPoint presentation - load existing if available, preserving user edits."""
        try:
            main_filename = "Introduction Module1.pptx"
            temp_filename = "Introduction Module1_TEMP.pptx"
            backup_filename = "Introduction Module1_BACKUP.pptx"
            
            # First try to load the main file
            if os.path.exists(main_filename):
                try:
                    print(f"Loading existing presentation: {main_filename}")
                    self.ppt = Presentation(main_filename)
                    existing_slides = len(self.ppt.slides)
                    print(f"Found {existing_slides} existing slides")
                    return
                except PermissionError:
                    print(f"Main file is open in PowerPoint, checking for temp file...")
                    
            # If main file doesn't exist or is locked, try temp file
            if os.path.exists(temp_filename):
                try:
                    print(f"Loading temp file: {temp_filename}")
                    self.ppt = Presentation(temp_filename)
                    return
                except:
                    print("Could not load temp file")
                    
            # Create new presentation if no existing files found
            print("Creating new presentation")
            self.ppt = Presentation()
            self.ppt.slide_width = Inches(13.33)  # Widescreen 16:9
            self.ppt.slide_height = Inches(7.5)
            
        except Exception as e:
            print(f"Error initializing PowerPoint: {e}")
            traceback.print_exc()
            self.ppt = Presentation()
            self.ppt.slide_width = Inches(13.33)
            self.ppt.slide_height = Inches(7.5)
    
    def detect_video_region(self):
        """Detect the video region on screen."""
        print("Detecting video region...")
        print("Please ensure your video is clearly visible on screen.")
        time.sleep(2)  # Give user time to position video
        
        # First try automatic detection
        self.video_region = self.auto_detect_video()
        
        # If automatic detection fails or results are poor, offer manual entry
        if not self.video_region:
            self.video_region = self.manual_region_entry()
            
        print(f"Final video region: {self.video_region}")
    
    def auto_detect_video(self):
        """Automatically detect video region using computer vision with focus on web video players."""
        try:
            # Take a full screenshot for processing
            screen = np.array(ImageGrab.grab())
            screen_size = (screen.shape[1], screen.shape[0])
            
            # Convert to different color spaces for better detection
            gray = cv2.cvtColor(screen, cv2.COLOR_RGB2GRAY)
            hsv = cv2.cvtColor(screen, cv2.COLOR_RGB2HSV)
            
            # 1. First try to detect Coursera video player specifically
            # Look for the video player controls bar which often has a distinct color
            lower_controls = np.array([0, 0, 40])  # Dark controls bar
            upper_controls = np.array([180, 30, 200])
            controls_mask = cv2.inRange(hsv, lower_controls, upper_controls)
            
            # 2. Look for the actual video content area which often has movement/contrast
            # Apply edge detection to find video boundaries
            edges = cv2.Canny(gray, 50, 150)
            
            # 3. Look for rectangles of the right size and aspect ratio
            # Combine masks for better detection
            combined_mask = cv2.bitwise_or(controls_mask, edges)
            
            # Dilate to connect nearby edges
            kernel = np.ones((5,5), np.uint8)
            dilated = cv2.dilate(combined_mask, kernel, iterations=2)
            
            # Find contours
            contours, _ = cv2.findContours(dilated, cv2.RETR_EXTERNAL, cv2.CHAIN_APPROX_SIMPLE)
            
            # Filter contours to find the video player
            best_rect = None
            max_score = 0
            
            for contour in contours:
                # Get approximate polygon
                epsilon = 0.02 * cv2.arcLength(contour, True)
                approx = cv2.approxPolyDP(contour, epsilon, True)
                
                # Check if it's rectangular (4-6 corners) and has minimum size
                if len(approx) >= 4 and len(approx) <= 8:
                    x, y, w, h = cv2.boundingRect(contour)
                    area = w * h
                    
                    # Skip tiny areas
                    if area < 10000:  # Minimum video area (e.g., 100x100)
                        continue
                        
                    # Skip if it's the entire screen
                    if w > screen_size[0] * 0.95 and h > screen_size[1] * 0.95:
                        continue
                    
                    # Calculate a score based on multiple factors
                    aspect_score = 0
                    # Prefer common video aspect ratios (16:9, 4:3, etc.)
                    aspect_ratio = w / h
                    if 1.7 < aspect_ratio < 1.8:  # 16:9
                        aspect_score = 1.0
                    elif 1.3 < aspect_ratio < 1.4:  # 4:3
                        aspect_score = 0.9
                    elif 1.0 < aspect_ratio < 2.0:  # Other reasonable video ratios
                        aspect_score = 0.7
                    else:
                        aspect_score = 0.3
                    
                    # Prefer rectangles that are not at the very edge of the screen
                    position_score = 0
                    if x > 10 and y > 10 and x + w < screen_size[0] - 10 and y + h < screen_size[1] - 10:
                        position_score = 1.0
                    else:
                        position_score = 0.5
                    
                    # Size score - prefer reasonably sized videos (not too small, not full screen)
                    size_ratio = area / (screen_size[0] * screen_size[1])
                    size_score = 0
                    if 0.1 < size_ratio < 0.7:  # Video takes up between 10% and 70% of screen
                        size_score = 1.0
                    elif 0.05 < size_ratio < 0.9:  # Bit smaller or larger
                        size_score = 0.7
                    else:
                        size_score = 0.3
                    
                    # Combine scores
                    final_score = (aspect_score * 0.4) + (position_score * 0.3) + (size_score * 0.3)
                    final_score *= area  # Larger areas with good scores are preferred
                    
                    if final_score > max_score:
                        max_score = final_score
                        best_rect = (x, y, x+w, y+h)
            
            # If we found a good rectangle, use it
            if best_rect:
                print(f"Video region auto-detected: {best_rect}")
                return best_rect
                
            # If no good rectangle found, try another approach: look for the play button/player controls
            # This is a fallback for when the video area itself isn't distinct enough
            lower_play = np.array([0, 0, 150])  # Whitish play button
            upper_play = np.array([180, 30, 255])
            play_mask = cv2.inRange(hsv, lower_play, upper_play)
            
            # Find contours of potential play buttons
            play_contours, _ = cv2.findContours(play_mask, cv2.RETR_EXTERNAL, cv2.CHAIN_APPROX_SIMPLE)
            
            # If we find a play button, assume the video is above it
            for contour in play_contours:
                x, y, w, h = cv2.boundingRect(contour)
                # A play button/controls bar is usually at the bottom of the video
                # Estimate video dimensions based on standard aspect ratios
                if w > 50 and h > 20:  # Minimum size for controls
                    video_width = w * 1.2  # Controls usually slightly narrower than video
                    video_height = video_width / 1.78  # Assume 16:9 aspect ratio
                    video_top = max(0, y - video_height)
                    video_left = max(0, x - (video_width - w) / 2)
                    video_rect = (
                        int(video_left), 
                        int(video_top), 
                        int(min(video_left + video_width, screen_size[0])), 
                        int(y + h)
                    )
                    return video_rect
                    
            # If all else fails, look for darker/video-like regions
            return None
            
        except Exception as e:
            print(f"Error in auto detection: {e}")
            return None
    
    def manual_region_entry(self):
        """Allow manual entry of video region coordinates."""
        print("\nAutomatic detection failed or produced suboptimal results.")
        print("Please enter the video region coordinates manually.")
        
        # Get screen dimensions for reference
        screen_size = ImageGrab.grab().size
        print(f"Your screen size is: {screen_size[0]}x{screen_size[1]}")
        
        try:
            # Create a simple dialog for coordinate entry
            root = tk.Tk()
            root.withdraw()  # Hide the main window
            
            # Show dialog with instructions
            message = (f"Enter video coordinates as: left,top,right,bottom\n"
                       f"Example: 100,100,700,500\n"
                       f"Your screen is {screen_size[0]}x{screen_size[1]}")
            
            coords_str = simpledialog.askstring("Enter Video Region", message)
            root.destroy()
            
            if coords_str:
                # Parse coordinates
                try:
                    left, top, right, bottom = map(int, coords_str.split(','))
                    return (left, top, right, bottom)
                except:
                    print("Invalid format. Using full screen.")
            
            # Fallback to full screen
            return (0, 0, screen_size[0], screen_size[1])
            
        except Exception as e:
            print(f"Error in manual entry: {e}")
            screen_size = ImageGrab.grab().size
            return (0, 0, screen_size[0], screen_size[1])
    
    def capture_loop(self):
        """Continuously capture screenshots while running."""
        try:
            while self.running:
                # Capture the video region
                screenshot = ImageGrab.grab(bbox=self.video_region)
                
                # Check if this image is similar to the last one
                if not self.is_duplicate(screenshot):
                    self.add_to_presentation(screenshot)
                    self.capture_count += 1
                    
                    # Save periodically as backup
                    if self.capture_count % self.save_interval == 0:
                        self.save_ppt(final=False)
                
                # Check if we should auto-stop
                elapsed = time.time() - self.start_time
                if elapsed > self.max_capture_time:
                    print(f"\nAUTO-TIMEOUT: Capture ran for {elapsed:.1f} seconds")
                    self.running = False
                    break
                    
                time.sleep(0.1)  # Capture frequency as needed
                
            # Final save when loop exits
            if not self.running:
                self.save_ppt(final=True)
                
        except Exception as e:
            print(f"Error in capture loop: {e}")
            traceback.print_exc()
            self.emergency_save()
    
    def is_duplicate(self, new_image):
        """Check if the new image is too similar to the last captured image."""
        if not self.last_image:
            return False
            
        try:
            # Convert images to numpy arrays for comparison
            new_array = np.array(new_image)
            last_array = np.array(self.last_image)
            
            # Ensure both images are the same size
            if new_array.shape != last_array.shape:
                return False
                
            # Calculate structural similarity index
            similarity = ssim(new_array, last_array, multichannel=True)
            
            # If similarity is above threshold, consider it a duplicate
            return similarity > self.similarity_threshold
            
        except Exception as e:
            print(f"Error in duplicate detection: {e}")
            return False
    
    def add_to_presentation(self, image):
        """Add the captured image to PowerPoint presentation."""
        try:
            # Save image temporarily
            temp_path = 'temp_screenshot.png'
            image.save(temp_path)
            
            # Add a new slide
            slide_layout = self.ppt.slide_layouts[5]  # Blank layout
            slide = self.ppt.slides.add_slide(slide_layout)
            
            # Add the image to completely fill the slide (no margins)
            left = Inches(0)
            top = Inches(0)
            width = self.ppt.slide_width
            height = self.ppt.slide_height
            
            slide.shapes.add_picture(temp_path, left, top, width, height)
            
            # Add to our collection
            self.images.append(image)
            self.last_image = image
            print(f"Added image {len(self.images)} to presentation (full-slide size)")
            
            # Save backup after first capture
            if len(self.images) == 1:
                self.save_ppt(final=False)
                
        except Exception as e:
            print(f"Error adding to presentation: {e}")
            traceback.print_exc()
    
    def save_ppt(self, final=True):
        """Save the PowerPoint presentation while preserving user edits."""
        if not self.ppt:
            return False
        
        try:
            main_filename = "Introduction Module1.pptx"
            temp_filename = "Introduction Module1_TEMP.pptx"
            
            if final:
                try:
                    # Try to save directly to main file
                    self.ppt.save(main_filename)
                    print(f"Saved to {main_filename}")
                    
                    # Clean up temp file if it exists
                    if os.path.exists(temp_filename):
                        try:
                            os.remove(temp_filename)
                        except:
                            pass
                        
                except PermissionError:
                    # If main file is open, save to temp file
                    self.ppt.save(temp_filename)
                    print(f"\nWARNING: PowerPoint file is open.")
                    print(f"Changes saved to: {temp_filename}")
                    print("Close PowerPoint and copy this temp file over the original")
                
            # Clean up temp screenshot if it exists
            if os.path.exists('temp_screenshot.png'):
                try:
                    os.remove('temp_screenshot.png')
                except:
                    pass
                
            return True
            
        except Exception as e:
            print(f"Error saving: {e}")
            traceback.print_exc()
            return False
        
    def emergency_save(self):
        """Last resort save attempt if something goes wrong."""
        try:
            print("\n!!! EMERGENCY SAVE ATTEMPT !!!")
            if self.ppt and self.images:
                try:
                    self.ppt.save("Introduction Module1_EMERGENCY.pptx")
                    print("Emergency save successful: Introduction Module1_EMERGENCY.pptx")
                except:
                    pass
                    
            # Brute force save of last image if everything else fails
            if self.last_image:
                try:
                    self.last_image.save("last_captured_image.png")
                    print("Last image saved: last_captured_image.png")
                except:
                    pass
        except:
            print("Emergency save failed completely")

# Global for keyboard listener
app = None
listener = None
last_key_time = 0
key_cooldown = 0.5  # Minimum seconds between key presses

def on_key_release(key):
    """Handle key release events."""
    global app, last_key_time
    try:
        # Check for key debouncing - prevent multiple rapid key triggers
        current_time = time.time()
        if current_time - last_key_time < key_cooldown:
            print(f"Ignoring rapid key press (wait {key_cooldown} seconds between presses)")
            return True
            
        # Update last key time
        last_key_time = current_time
        
        # Check for END key
        if key == keyboard.Key.end:
            print("END key pressed - taking a screenshot")
            app.toggle_capture()
        # Alternative key: F12
        elif key == keyboard.Key.f12:
            print("F12 key pressed - taking a screenshot")
            app.toggle_capture()
        # Alternative key: pressing 'p'
        elif hasattr(key, 'char') and key.char == 'p':
            print("'p' key pressed - taking a screenshot")
            app.toggle_capture()
        # Exit on ESC
        elif key == keyboard.Key.esc:
            print("ESC key pressed - exiting program")
            # Stop the listener
            return False
    except Exception as e:
        print(f"Error handling key: {e}")
        if app:
            app.emergency_save()
        
    return True

def emergency_exit_handler():
    global app
    if app and app.running:
        try:
            print("\n!!! EMERGENCY EXIT TRIGGERED !!!")
            app.running = False
            app.emergency_save()
        except:
            pass

if __name__ == "__main__":
    try:
        app = VideoToPPT()
        
        # Register emergency exit handler
        import atexit
        atexit.register(emergency_exit_handler)
        
        print("Video to PowerPoint Screenshot Tool")
        print("===================================")
        print("Press END, F12, or 'p' key to take a screenshot")
        print("Each key press captures one screenshot")
        print("Press ESC to exit the program")
        
        # Start listening for key presses
        with keyboard.Listener(on_release=on_key_release) as listener:
            listener.join()
            
    except Exception as e:
        print(f"CRITICAL ERROR: {e}")
        traceback.print_exc()
        if app:
            app.emergency_save()
        
        print("\nPress ENTER to exit...")
        input() 